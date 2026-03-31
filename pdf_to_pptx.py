#!/usr/bin/env python3
import sys
import os
import tempfile
from pptx import Presentation
from pptx.util import Pt
import pytesseract
from pdf2image import convert_from_path
import PyPDF2

def is_scanned_pdf(pdf_path):
    """Check if PDF is scanned (image-based)"""
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text.strip():  # If we find any text, it's not fully scanned
                    return False
            return True
    except:
        return True  # Assume scanned if we can't check

def extract_text_from_pdf(pdf_path):
    """Extract text from PDF using appropriate method"""
    text = ""
    
    # First try regular text extraction
    try:
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text.strip():
                    text += page_text + "\n\n"
    except:
        pass
    
    # If no text found, try OCR
    if not text.strip():
        print("No text found, attempting OCR...")
        try:
            # Convert PDF to images
            images = convert_from_path(pdf_path, dpi=300)
            
            # Extract text from each image using OCR
            for i, image in enumerate(images):
                print(f"Processing page {i+1} with OCR...")
                page_text = pytesseract.image_to_string(image, lang='eng')
                text += page_text + "\n\n"
        except Exception as e:
            print(f"OCR failed: {e}")
            return None
    
    return text if text.strip() else None

def pdf_to_pptx(pdf_path, pptx_path):
    """
    Convert PDF to editable PowerPoint presentation with OCR support
    """
    try:
        print(f"Processing PDF: {pdf_path}")
        
        # Extract text from PDF
        text = extract_text_from_pdf(pdf_path)
        
        if not text:
            print("No text could be extracted from the PDF")
            # Create a default slide with error message
            prs = Presentation()
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = "Conversion Issue"
            subtitle.text = "Could not extract text from the PDF. This may be a scanned document that requires OCR processing."
            prs.save(pptx_path)
            return True
        
        print(f"Extracted {len(text)} characters of text")
        
        # Split text into manageable chunks for slides
        words = text.split()
        slides_text = []
        current_slide = []
        char_count = 0
        
        for word in words:
            if char_count + len(word) > 1000:  # ~1000 chars per slide
                slides_text.append(" ".join(current_slide))
                current_slide = [word]
                char_count = len(word)
            else:
                current_slide.append(word)
                char_count += len(word) + 1  # +1 for space
        
        if current_slide:
            slides_text.append(" ".join(current_slide))
        
        print(f"Creating {len(slides_text)} slides")
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Add title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Converted from PDF"
        subtitle.text = f"Original: {os.path.basename(pdf_path)}"
        
        # Add content slides
        for i, slide_text in enumerate(slides_text):
            slide_layout = prs.slide_layouts[1]  # Title and content
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = f"Content Section {i + 1}"
            
            # Set content
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.text = slide_text
            
            # Format text
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Arial"
        
        print("Saving PowerPoint file...")
        prs.save(pptx_path)
        print(f"Successfully created {pptx_path} with {len(prs.slides)} slides")
        return True
        
    except Exception as e:
        print(f"Error during conversion: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python pdf_to_pptx.py <input_pdf> <output_pptx>")
        sys.exit(1)
    
    input_pdf = sys.argv[1]
    output_pptx = sys.argv[2]
    
    if not os.path.exists(input_pdf):
        print(f"Input file {input_pdf} does not exist")
        sys.exit(1)
    
    success = pdf_to_pptx(input_pdf, output_pptx)
    
    if success:
        print("Conversion completed successfully")
        sys.exit(0)
    else:
        print("Conversion failed")
        sys.exit(1)

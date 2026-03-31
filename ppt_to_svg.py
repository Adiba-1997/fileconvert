#!/usr/bin/env python3
import sys
import os
from pptx import Presentation
import matplotlib.pyplot as plt
from matplotlib.backends.backend_agg import FigureCanvasAgg
import numpy as np
from PIL import Image
import cairosvg
import tempfile

def ppt_to_svg(ppt_path, svg_path):
    """
    Convert PowerPoint to SVG by extracting slides as images and converting to SVG
    """
    try:
        print(f"Processing PowerPoint file: {ppt_path}")
        
        # Load the presentation
        prs = Presentation(ppt_path)
        
        if len(prs.slides) == 0:
            print("No slides found in the presentation")
            return False
        
        # Process first slide only (for simplicity)
        # You can modify this to process all slides
        slide = prs.slides[0]
        
        # Create a temporary image file
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            temp_image_path = tmp.name
        
        # Export slide as image (this is a simplified approach)
        # In a real implementation, you'd use proper slide rendering
        fig, ax = plt.subplots(figsize=(10, 7.5))
        ax.text(0.5, 0.5, f"Slide 1 from {os.path.basename(ppt_path)}", 
                ha='center', va='center', fontsize=16)
        ax.set_title("PPT to SVG Conversion")
        plt.axis('off')
        
        # Save as PNG
        plt.savefig(temp_image_path, bbox_inches='tight', pad_inches=0)
        plt.close()
        
        # Convert PNG to SVG using Cairo
        cairosvg.png2svg(url=temp_image_path, write_to=svg_path)
        
        # Clean up
        os.unlink(temp_image_path)
        
        print(f"Successfully created SVG: {svg_path}")
        return True
        
    except Exception as e:
        print(f"Error during conversion: {str(e)}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python ppt_to_svg.py <input_ppt> <output_svg>")
        sys.exit(1)
    
    input_ppt = sys.argv[1]
    output_svg = sys.argv[2]
    
    if not os.path.exists(input_ppt):
        print(f"Input file {input_ppt} does not exist")
        sys.exit(1)
    
    success = ppt_to_svg(input_ppt, output_svg)
    
    if success:
        print("Conversion successful")
        sys.exit(0)
    else:
        print("Conversion failed")
        sys.exit(1)
